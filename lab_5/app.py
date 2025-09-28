from flask import Flask, render_template, request, jsonify, send_file, flash
import pandas as pd
import json
import xml.etree.ElementTree as ET
from io import StringIO, BytesIO
import base64
from pathlib import Path
import openpyxl
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from werkzeug.utils import secure_filename
import os
import tempfile

# Extended format support imports
import PyPDF2
from docx import Document
try:
    import pymupdf as fitz
except ImportError:
    fitz = None
import pdfplumber
from pptx import Presentation
from PIL import Image
import markdown
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# Imports for my lab
import logging
from logging.handlers import RotatingFileHandler
import sys
from datetime import datetime

app = Flask(__name__)
app.secret_key = 'your_secret_key_here'

# Configure upload settings
UPLOAD_FOLDER = 'temp_uploads'
ALLOWED_EXTENSIONS = {
    'csv', 'json', 'xml', 'txt', 'xlsx', 'pdf', 'docx', 'pptx', 
    'jpg', 'jpeg', 'png', 'bmp', 'gif', 'md', 'html'
}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

# Create upload directory if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


# Настройка логирования
def setup_logging():
    """Настройка системы логирования"""
    # Создаем папку для логов если ее нет
    log_dir = 'logs'
    os.makedirs(log_dir, exist_ok=True)

    # Форматтер для логов
    formatter = logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(message)s - [%(filename)s:%(lineno)d]'
    )

    # Логгер приложения
    logger = logging.getLogger('file_converter')
    logger.setLevel(logging.DEBUG)

    # Файловый обработчик с ротацией
    file_handler = RotatingFileHandler(
        f'{log_dir}/converter.log',
        maxBytes=10 * 1024 * 1024,  # 10MB
        backupCount=5,
        encoding='utf-8'
    )
    file_handler.setLevel(logging.DEBUG)
    file_handler.setFormatter(formatter)

    # Консольный обработчик
    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setLevel(logging.INFO)
    console_handler.setFormatter(formatter)

    # Добавляем обработчики
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)

    return logger


# Инициализация логгера
logger = setup_logging()


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def convert_csv_to_json(csv_content):
    """Convert CSV to JSON"""
    logger.debug("Конвертация CSV в JSON")
    try:
        df = pd.read_csv(StringIO(csv_content))
        result = df.to_json(orient='records', indent=2)
        logger.debug(f"✅ CSV→JSON: обработано {len(df)} строк")
        return result
    except Exception as e:
        logger.error(f"❌ Ошибка конвертации CSV в JSON: {str(e)}")
        raise


def convert_csv_to_xml(csv_content):
    """Convert CSV to XML"""
    df = pd.read_csv(StringIO(csv_content))
    xml_data = []
    xml_data.append('<?xml version="1.0" encoding="UTF-8"?>')
    xml_data.append('<root>')
    
    for _, row in df.iterrows():
        xml_data.append('  <record>')
        for col in df.columns:
            xml_data.append(f'    <{col}>{row[col]}</{col}>')
        xml_data.append('  </record>')
    
    xml_data.append('</root>')
    return '\n'.join(xml_data)


def convert_json_to_csv(json_content):
    """Convert JSON to CSV"""
    data = json.loads(json_content)
    if isinstance(data, list):
        df = pd.DataFrame(data)
    else:
        df = pd.DataFrame([data])
    return df.to_csv(index=False)


def convert_json_to_xml(json_content):
    """Convert JSON to XML"""
    data = json.loads(json_content)
    
    def dict_to_xml(d, root_name="root"):
        root = ET.Element(root_name)
        
        def add_to_xml(parent, key, value):
            if isinstance(value, dict):
                child = ET.SubElement(parent, key)
                for k, v in value.items():
                    add_to_xml(child, k, v)
            elif isinstance(value, list):
                for item in value:
                    if isinstance(item, dict):
                        child = ET.SubElement(parent, key)
                        for k, v in item.items():
                            add_to_xml(child, k, v)
                    else:
                        child = ET.SubElement(parent, key)
                        child.text = str(item)
            else:
                child = ET.SubElement(parent, key)
                child.text = str(value)
        
        if isinstance(data, list):
            for item in data:
                add_to_xml(root, "item", item)
        else:
            for key, value in data.items():
                add_to_xml(root, key, value)
        
        return ET.tostring(root, encoding='unicode')
    
    return dict_to_xml(data)


def convert_csv_to_excel(csv_content):
    """Convert CSV to Excel"""
    df = pd.read_csv(StringIO(csv_content))
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Sheet1')
    return output.getvalue()


def convert_to_txt(content, source_format):
    """Convert various formats to plain text"""
    if source_format == "csv":
        df = pd.read_csv(StringIO(content))
        return df.to_string(index=False)
    elif source_format == "json":
        data = json.loads(content)
        return json.dumps(data, indent=2)
    elif source_format == "xml":
        root = ET.fromstring(content)
        return ET.tostring(root, encoding='unicode')
    else:
        return content


def convert_to_pdf(content, source_format):
    """Convert content to PDF"""
    text_content = convert_to_txt(content, source_format)
    
    buffer = BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    lines = text_content.split('\n')
    y_position = height - 50
    line_height = 12
    
    for line in lines:
        if y_position < 50:
            p.showPage()
            y_position = height - 50
        
        if len(line) > 100:
            line = line[:97] + "..."
        
        p.drawString(50, y_position, line)
        y_position -= line_height
    
    p.save()
    return buffer.getvalue()


# Extended conversion functions
def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    logger.debug("🔍 Извлечение текста из PDF")
    text = ""
    try:
        with pdfplumber.open(pdf_file) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                    logger.debug(f"📄 Страница {page_num}: извлечено {len(page_text)} символов")
                else:
                    logger.warning(f"⚠️ Страница {page_num}: текст не извлечен")

        logger.info(f"✅ Извлечено {len(text)} символов из PDF")
        return text
    except Exception as e:
        try:
            # Fallback to PyMuPDF if available
            logger.error(f"❌ Ошибка извлечения текста из PDF: {str(e)}")
            if fitz:
                doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
                for page in doc:
                    text += page.get_text() + "\n"
                doc.close()
            else:
                raise Exception("PDF text extraction failed and PyMuPDF not available")
        except Exception as fallback_error:
            # Last resort: return error message
            text = f"Error extracting text from PDF: {str(e)}\nFallback error: {str(fallback_error)}"
    return text


def extract_text_from_docx(docx_file):
    """Extract text from Word document"""
    doc = Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text


def extract_text_from_pptx(pptx_file):
    """Extract text from PowerPoint presentation with enhanced content extraction"""
    prs = Presentation(pptx_file)
    full_text = ""
    slide_number = 0
    
    for slide in prs.slides:
        slide_number += 1
        slide_text = f"\n{'='*50}\nSLIDE {slide_number}\n{'='*50}\n\n"
        
        # Extract text from all shapes
        for shape in slide.shapes:
            try:
                # Text in text boxes and shapes
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n\n"
                
                # Text in tables
                elif shape.has_table:
                    table = shape.table
                    slide_text += "[TABLE CONTENT]\n"
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_text += " | ".join(row_text) + "\n"
                    slide_text += "\n"
                
                # Text in grouped shapes
                elif hasattr(shape, 'shapes'):
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text") and sub_shape.text.strip():
                            slide_text += sub_shape.text.strip() + "\n"
                
            except Exception as e:
                # Continue processing other shapes if one fails
                continue
        
        # Extract notes if present
        try:
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text += f"\n[SLIDE NOTES]\n{notes_text}\n\n"
        except:
            pass
        
        # Only add slide if it has content
        if slide_text.strip() != f"\n{'='*50}\nSLIDE {slide_number}\n{'='*50}\n\n":
            full_text += slide_text
        else:
            # Add placeholder for slides with no extractable text
            full_text += f"\n{'='*50}\nSLIDE {slide_number}\n{'='*50}\n\n[Slide contains visual content that cannot be extracted as text]\n\n"
    
    return full_text if full_text.strip() else "No extractable text content found in presentation."


def convert_text_to_docx(text_content):
    """Convert text to Word document"""
    doc = Document()
    
    # Split text into paragraphs
    paragraphs = text_content.split('\n\n')
    
    for para_text in paragraphs:
        if para_text.strip():
            doc.add_paragraph(para_text.strip())
    
    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def convert_text_to_html(text_content):
    """Convert text to HTML"""
    html = f"""<!DOCTYPE html>
<html>
<head>
    <title>Converted Document</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1, h2, h3 {{ color: #333; }}
        p {{ margin-bottom: 1em; }}
    </style>
</head>
<body>
    <h1>Converted Document</h1>
    <div>
        {text_content.replace(chr(10), '<br>')}
    </div>
</body>
</html>"""
    return html


def convert_markdown_to_html(md_content):
    """Convert Markdown to HTML"""
    html_content = markdown.markdown(md_content)
    return f"""<!DOCTYPE html>
<html>
<head>
    <title>Converted from Markdown</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1, h2, h3 {{ color: #333; }}
        code {{ background: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
        pre {{ background: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }}
    </style>
</head>
<body>
    {html_content}
</body>
</html>"""


def convert_image_format(image_file, target_format):
    """Convert image between different formats"""
    logger.info(f"🖼️ Конвертация изображения в {target_format}")
    try:
        img = Image.open(image_file)
        original_format = img.format
        logger.debug(f"Исходный формат: {original_format}, размер: {img.size}")

        # ... код конвертации ...

        logger.info(f"✅ Изображение сконвертировано: {original_format} → {target_format}")
        return buffer.getvalue()
    except Exception as e:
        logger.error(f"❌ Ошибка конвертации изображения: {str(e)}")
        raise


def create_pdf_from_text(text_content):
    """Create a formatted PDF from text using ReportLab with better PowerPoint handling"""
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.platypus import PageBreak
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Create custom styles
    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontSize=16,
        spaceAfter=20,
        alignment=TA_CENTER,
        textColor='#0066cc'
    )
    
    # Check if this looks like PowerPoint content (has slide markers)
    if 'SLIDE ' in text_content and '=' in text_content:
        # Handle PowerPoint-style content
        slides = text_content.split('\n' + '='*50)
        
        for i, slide_content in enumerate(slides):
            if slide_content.strip():
                # Clean up slide content
                slide_content = slide_content.strip()
                
                # Extract slide title/number
                lines = slide_content.split('\n')
                slide_title = None
                content_lines = []
                
                for line in lines:
                    if line.startswith('SLIDE '):
                        slide_title = line
                    elif line.strip() and not line.startswith('='):
                        content_lines.append(line)
                
                # Add slide title
                if slide_title:
                    if i > 0:  # Add page break before new slides (except first)
                        story.append(PageBreak())
                    title_para = Paragraph(slide_title, slide_title_style)
                    story.append(title_para)
                    story.append(Spacer(1, 20))
                
                # Add slide content
                for line in content_lines:
                    if line.strip():
                        if line.startswith('[TABLE CONTENT]'):
                            # Style table content differently
                            para = Paragraph(line, styles['Heading3'])
                        elif line.startswith('[SLIDE NOTES]'):
                            # Style notes differently
                            para = Paragraph(line, styles['Heading3'])
                        elif '|' in line and len(line.split('|')) > 1:
                            # This looks like table data
                            para = Paragraph(line.replace('|', ' | '), styles['Code'])
                        else:
                            para = Paragraph(line, styles['Normal'])
                        
                        story.append(para)
                        story.append(Spacer(1, 8))
    else:
        # Handle regular text content
        paragraphs = text_content.split('\n\n')
        
        for para_text in paragraphs:
            if para_text.strip():
                # Split very long paragraphs
                if len(para_text) > 1000:
                    sentences = para_text.split('. ')
                    current_para = ""
                    
                    for sentence in sentences:
                        if len(current_para + sentence) > 800:
                            if current_para:
                                para = Paragraph(current_para.strip(), styles['Normal'])
                                story.append(para)
                                story.append(Spacer(1, 12))
                            current_para = sentence + '. '
                        else:
                            current_para += sentence + '. '
                    
                    if current_para.strip():
                        para = Paragraph(current_para.strip(), styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 12))
                else:
                    para = Paragraph(para_text.strip(), styles['Normal'])
                    story.append(para)
                    story.append(Spacer(1, 12))
    
    # Build the PDF
    doc.build(story)
    return buffer.getvalue()


def perform_conversion(file_content, input_format, target_format, file_obj=None):
    """Perform file conversion based on formats"""
    logger.debug(f"🛠️ Вызов perform_conversion: {input_format} -> {target_format}")

    try:
        # Логируем попытку конвертации
        logger.info(f"🔧 Конвертация: {input_format.upper()} → {target_format.upper()}")

        # Сохраняем результат конвертации
        converted_content = None

        # Data format conversions (CSV, JSON, XML)
        if input_format == 'csv':
            if target_format == 'json':
                logger.debug("Конвертация CSV → JSON")
                converted_content = convert_csv_to_json(file_content)
            elif target_format == 'xml':
                logger.debug("Конвертация CSV → XML")
                converted_content = convert_csv_to_xml(file_content)
            elif target_format == 'xlsx':
                logger.debug("Конвертация CSV → Excel")
                converted_content = convert_csv_to_excel(file_content)
            elif target_format == 'txt':
                logger.debug("Конвертация CSV → TXT")
                converted_content = convert_to_txt(file_content, 'csv')
            elif target_format == 'pdf':
                logger.debug("Конвертация CSV → PDF")
                converted_content = create_pdf_from_text(file_content)
            elif target_format == 'docx':
                logger.debug("Конвертация CSV → DOCX")
                converted_content = convert_text_to_docx(file_content)
            elif target_format == 'html':
                logger.debug("Конвертация CSV → HTML")
                converted_content = convert_text_to_html(file_content)

        elif input_format == 'json':
            if target_format == 'csv':
                logger.debug("Конвертация JSON → CSV")
                converted_content = convert_json_to_csv(file_content)
            elif target_format == 'xml':
                logger.debug("Конвертация JSON → XML")
                converted_content = convert_json_to_xml(file_content)
            elif target_format == 'xlsx':
                logger.debug("Конвертация JSON → Excel")
                csv_content = convert_json_to_csv(file_content)
                converted_content = convert_csv_to_excel(csv_content)
            elif target_format == 'txt':
                logger.debug("Конвертация JSON → TXT")
                converted_content = convert_to_txt(file_content, 'json')
            elif target_format == 'pdf':
                logger.debug("Конвертация JSON → PDF")
                text_content = convert_to_txt(file_content, 'json')
                converted_content = create_pdf_from_text(text_content)
            elif target_format == 'docx':
                logger.debug("Конвертация JSON → DOCX")
                text_content = convert_to_txt(file_content, 'json')
                converted_content = convert_text_to_docx(text_content)
            elif target_format == 'html':
                logger.debug("Конвертация JSON → HTML")
                text_content = convert_to_txt(file_content, 'json')
                converted_content = convert_text_to_html(text_content)

        elif input_format in ['xml', 'txt']:
            if target_format == 'txt':
                logger.debug(f"Конвертация {input_format.upper()} → TXT")
                converted_content = convert_to_txt(file_content, input_format)
            elif target_format == 'pdf':
                logger.debug(f"Конвертация {input_format.upper()} → PDF")
                converted_content = create_pdf_from_text(file_content)
            elif target_format == 'docx':
                logger.debug(f"Конвертация {input_format.upper()} → DOCX")
                converted_content = convert_text_to_docx(file_content)
            elif target_format == 'html':
                logger.debug(f"Конвертация {input_format.upper()} → HTML")
                converted_content = convert_text_to_html(file_content)
            elif target_format == 'json':
                logger.debug(f"Конвертация {input_format.upper()} → JSON")
                # Try to convert text to JSON if it looks like structured data
                try:
                    data = json.loads(file_content)
                    converted_content = json.dumps(data, indent=2)
                except:
                    # Convert to simple JSON structure
                    converted_content = json.dumps({"content": file_content}, indent=2)

        # Document format conversions (PDF, DOCX, PPTX)
        elif input_format == 'pdf':
            logger.debug("Обработка PDF файла")
            if file_obj:
                text_content = extract_text_from_pdf(file_obj)
            else:
                text_content = file_content

            logger.debug(f"Извлечено {len(text_content)} символов из PDF")

            if target_format == 'txt':
                logger.debug("Конвертация PDF → TXT")
                converted_content = text_content
            elif target_format == 'docx':
                logger.debug("Конвертация PDF → DOCX")
                converted_content = convert_text_to_docx(text_content)
            elif target_format == 'html':
                logger.debug("Конвертация PDF → HTML")
                converted_content = convert_text_to_html(text_content)
            elif target_format == 'json':
                logger.debug("Конвертация PDF → JSON")
                converted_content = json.dumps({"extracted_text": text_content}, indent=2)
            elif target_format == 'csv':
                logger.debug("Конвертация PDF → CSV")
                # Create a simple CSV with the extracted text
                converted_content = f"extracted_text\n{text_content}"

        elif input_format == 'docx':
            logger.debug("Обработка DOCX файла")
            if file_obj:
                text_content = extract_text_from_docx(file_obj)
            else:
                text_content = file_content

            logger.debug(f"Извлечено {len(text_content)} символов из DOCX")

            if target_format == 'txt':
                logger.debug("Конвертация DOCX → TXT")
                converted_content = text_content
            elif target_format == 'pdf':
                logger.debug("Конвертация DOCX → PDF")
                converted_content = create_pdf_from_text(text_content)
            elif target_format == 'html':
                logger.debug("Конвертация DOCX → HTML")
                converted_content = convert_text_to_html(text_content)
            elif target_format == 'json':
                logger.debug("Конвертация DOCX → JSON")
                converted_content = json.dumps({"extracted_text": text_content}, indent=2)
            elif target_format == 'csv':
                logger.debug("Конвертация DOCX → CSV")
                converted_content = f"extracted_text\n{text_content}"

        elif input_format == 'pptx':
            logger.debug("Обработка PPTX файла")
            if file_obj:
                text_content = extract_text_from_pptx(file_obj)
            else:
                text_content = file_content

            logger.debug(f"Извлечено {len(text_content)} символов из PPTX")

            if target_format == 'txt':
                logger.debug("Конвертация PPTX → TXT")
                converted_content = text_content
            elif target_format == 'pdf':
                logger.debug("Конвертация PPTX → PDF")
                converted_content = create_pdf_from_text(text_content)
            elif target_format == 'docx':
                logger.debug("Конвертация PPTX → DOCX")
                converted_content = convert_text_to_docx(text_content)
            elif target_format == 'html':
                logger.debug("Конвертация PPTX → HTML")
                converted_content = convert_text_to_html(text_content)
            elif target_format == 'json':
                logger.debug("Конвертация PPTX → JSON")
                converted_content = json.dumps({"extracted_text": text_content}, indent=2)

        # Markdown conversions
        elif input_format == 'md':
            logger.debug("Обработка Markdown файла")
            if target_format == 'html':
                logger.debug("Конвертация MD → HTML")
                converted_content = convert_markdown_to_html(file_content)
            elif target_format == 'txt':
                logger.debug("Конвертация MD → TXT")
                converted_content = file_content
            elif target_format == 'pdf':
                logger.debug("Конвертация MD → PDF")
                html_content = convert_markdown_to_html(file_content)
                # For now, convert HTML content as text to PDF
                converted_content = create_pdf_from_text(file_content)
            elif target_format == 'docx':
                logger.debug("Конвертация MD → DOCX")
                converted_content = convert_text_to_docx(file_content)

        # HTML conversions
        elif input_format == 'html':
            logger.debug("Обработка HTML файла")
            if target_format == 'txt':
                logger.debug("Конвертация HTML → TXT")
                # Strip HTML tags for text conversion
                import re
                clean_text = re.sub('<[^<]+?>', '', file_content)
                converted_content = clean_text
            elif target_format == 'pdf':
                logger.debug("Конвертация HTML → PDF")
                clean_text = re.sub('<[^<]+?>', '', file_content)
                converted_content = create_pdf_from_text(clean_text)
            elif target_format == 'docx':
                logger.debug("Конвертация HTML → DOCX")
                clean_text = re.sub('<[^<]+?>', '', file_content)
                converted_content = convert_text_to_docx(clean_text)

        # Image format conversions
        elif input_format in ['jpg', 'jpeg', 'png', 'bmp', 'gif']:
            logger.debug(f"Конвертация изображения {input_format} → {target_format}")
            if target_format in ['jpg', 'jpeg', 'png', 'bmp', 'gif']:
                if file_obj:
                    converted_content = convert_image_format(file_obj, target_format)
                else:
                    raise ValueError("Image file object required for image conversion")

        else:
            logger.error(f"❌ Неподдерживаемый входной формат: {input_format}")
            raise ValueError(f"Unsupported input format: {input_format}")

        # Проверяем что конвертация прошла успешно
        if converted_content is None:
            logger.error(f"❌ Конвертация не выполнена: {input_format} -> {target_format}")
            raise ValueError(f"Conversion from {input_format} to {target_format} is not supported")

        logger.info(f"✅ Успешная конвертация: {input_format} -> {target_format}")
        return converted_content

    except Exception as e:
        logger.error(f"❌ Ошибка конвертации {input_format} -> {target_format}: {str(e)}", exc_info=True)
        raise ValueError(f"Conversion failed: {str(e)}")


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/convert', methods=['POST'])
def convert_file():
    logger.info("📥 Получен запрос на конвертацию файла")

    try:
        if 'file' not in request.files:
            logger.warning("❌ В запросе отсутствует файл")
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        target_format = request.form.get('target_format')

        if file.filename == '':
            logger.warning("❌ Не выбрано имя файла")
            return jsonify({'error': 'No file selected'}), 400

        if not allowed_file(file.filename):
            logger.warning(f"🚫 Неподдерживаемый тип файла: {file.filename}")
            return jsonify({'error': 'File type not supported'}), 400

        if not target_format:
            logger.warning("❌ Не указан целевой формат")
            return jsonify({'error': 'Target format not specified'}), 400

        # Логируем информацию о файле
        input_format = Path(file.filename).suffix.lower()[1:]

        # Читаем размер файла для логирования
        file.seek(0, 2)  # Перемещаемся в конец файла
        file_size = file.tell()  # Получаем размер
        file.seek(0)  # Возвращаем указатель в начало

        logger.info(f"📄 Конвертация файла: {file.filename} "
                    f"(размер: {file_size} байт, "
                    f"из: {input_format} в: {target_format})")

        # Оригинальный код обработки файла
        file_content = None
        file_obj = None

        # Reset file pointer
        file.seek(0)

        # Binary formats that need special handling
        binary_formats = ['pdf', 'docx', 'pptx', 'xlsx', 'jpg', 'jpeg', 'png', 'bmp', 'gif']

        if input_format in binary_formats:
            file_obj = file
            if input_format == 'xlsx':
                # Special handling for Excel files
                logger.debug("🔍 Обработка Excel файла")
                df = pd.read_excel(file)
                file_content = df.to_csv(index=False)
                input_format = 'csv'
                file_obj = None
                logger.debug(f"📊 Excel конвертирован в CSV: {len(df)} строк")
            else:
                file_content = ""  # Will be processed using file_obj
                logger.debug(f"🔧 Бинарный файл будет обработан через file_obj: {input_format}")
        else:
            # Text formats
            try:
                file_content = file.read().decode('utf-8')
                logger.debug(f"📝 Текстовый файл прочитан: {len(file_content)} символов")
            except UnicodeDecodeError:
                logger.error("❌ Ошибка декодирования файла - не текстовый формат")
                return jsonify({'error': 'Unable to decode file. Please ensure it\'s a text file.'}), 400

        # Выполняем конвертацию
        logger.debug(f"🔄 Начало конвертации {input_format} -> {target_format}")
        converted_content = perform_conversion(file_content, input_format, target_format, file_obj)
        logger.debug("✅ Конвертация завершена успешно")

        # Генерируем имя выходного файла
        original_name = Path(file.filename).stem
        output_filename = f"{original_name}.{target_format}"

        # Сохраняем сконвертированный файл
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{target_format}")

        # Определяем бинарные форматы
        binary_target_formats = ['xlsx', 'pdf', 'docx', 'pptx', 'jpg', 'jpeg', 'png', 'bmp', 'gif']

        # Создаем preview для текстовых форматов
        preview = None
        if target_format not in binary_target_formats and isinstance(converted_content, str):
            preview = converted_content[:500] + "..." if len(converted_content) > 500 else converted_content
            logger.debug(f"👀 Создан preview: {len(preview)} символов")

        if target_format in binary_target_formats:
            temp_file.write(converted_content)
            logger.debug(f"💾 Сохранен бинарный файл: {temp_file.name}")
        else:
            if isinstance(converted_content, str):
                temp_file.write(converted_content.encode('utf-8'))
                logger.debug(f"💾 Сохранен текстовый файл: {temp_file.name}, размер: {len(converted_content)} символов")
            else:
                temp_file.write(converted_content)
                logger.debug(f"💾 Сохранен файл: {temp_file.name}, бинарный размер: {len(converted_content)} байт")

        temp_file.close()

        # Генерируем уникальный ID
        import uuid
        file_id = str(uuid.uuid4())
        temp_files[file_id] = temp_file.name

        logger.info(f"✅ Конвертация завершена успешно. ID файла: {file_id}, "
                    f"результирующее имя: {output_filename}")

        return jsonify({
            'success': True,
            'filename': output_filename,
            'file_id': file_id,
            'preview': preview
        })

    except Exception as e:
        logger.error(f"💥 Ошибка при конвертации: {str(e)}", exc_info=True)
        return jsonify({'error': str(e)}), 500


# Store temporary files with unique IDs
temp_files = {}


@app.route('/download/<file_id>/<filename>')
def download_file(file_id, filename):
    logger.info(f"📥 Запрос на скачивание файла: {filename} (ID: {file_id})")

    try:
        if file_id not in temp_files:
            logger.warning(f"🚫 Файл не найден или устарел: {file_id}")
            return jsonify({'error': 'File not found or expired'}), 404

        file_path = temp_files[file_id]

        if not os.path.exists(file_path):
            logger.error(f"🚫 Физический файл отсутствует: {file_path}")
            return jsonify({'error': f'File not found: {file_path}'}), 404

        file_size = os.path.getsize(file_path)
        logger.info(f"📤 Отправка файла: {filename} (размер: {file_size} байт)")

        response = send_file(file_path, as_attachment=True, download_name=filename)

        @response.call_on_close
        def cleanup():
            try:
                os.unlink(file_path)
                temp_files.pop(file_id, None)
                logger.debug(f"🧹 Очистка временного файла: {file_path}")
            except Exception as e:
                logger.error(f"❌ Ошибка при очистке файла {file_path}: {str(e)}")

        return response

    except Exception as e:
        logger.error(f"💥 Ошибка скачивания файла {filename}: {str(e)}", exc_info=True)
        return jsonify({'error': f'Download failed: {str(e)}'}), 500


if __name__ == '__main__':
    logger.info("🚀 Запуск Flask приложения File Converter")
    logger.info("📍 Хост: 0.0.0.0, Порт: 5000")
    logger.info("🔧 Режим отладки: ВКЛ" if app.debug else "🔧 Режим отладки: ВЫКЛ")

    try:
        app.run(debug=True, host='0.0.0.0', port=5000)
    except Exception as e:
        logger.critical(f"💥 Критическая ошибка при запуске приложения: {str(e)}")
        raise
